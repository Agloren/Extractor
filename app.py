import streamlit as st
import fitz
import io
import json
import re
import tempfile
from anthropic import Anthropic
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from moviepy import VideoFileClip

# ─────────────────────────────────────
# CONFIG
# ─────────────────────────────────────
st.set_page_config(page_title="AI Study Buddy", layout="wide")

if "ANTHROPIC_API_KEY" not in st.secrets:
    st.error("Falta ANTHROPIC_API_KEY en secrets.")
    st.stop()

client = Anthropic(api_key=st.secrets["ANTHROPIC_API_KEY"])

# ─────────────────────────────────────
# HELPERS
# ─────────────────────────────────────
def hex_to_rgb(hex_color):
    hex_color = hex_color.lstrip("#")
    return RGBColor(
        int(hex_color[0:2], 16),
        int(hex_color[2:4], 16),
        int(hex_color[4:6], 16)
    )

def smart_truncate(text, max_chars=15000):
    return text[:max_chars]

def safe_json_parse(text):
    text = re.sub(r"```json|```", "", text).strip()
    return json.loads(text)

# ─────────────────────────────────────
# EXTRACTORES
# ─────────────────────────────────────
def extract_pdf(file_bytes):
    doc = fitz.open(stream=file_bytes, filetype="pdf")
    return "\n".join([p.get_text() for p in doc])

def extract_docx(file_bytes):
    import docx
    doc = docx.Document(io.BytesIO(file_bytes))
    return "\n".join([p.text for p in doc.paragraphs])

def extract_txt(file_bytes):
    return file_bytes.decode("utf-8", errors="ignore")

def extract_pptx(file_bytes):
    prs = Presentation(io.BytesIO(file_bytes))
    text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)
    return "\n".join(text)

def transcribe_audio(file_bytes, filename):
    response = client.messages.create(
        model="claude-sonnet-4-5-20250929",
        max_tokens=4000,
        messages=[{
            "role": "user",
            "content": "Transcribe fielmente el siguiente audio."
        }],
        attachments=[{
            "file_name": filename,
            "mime_type": "audio/mpeg",
            "data": file_bytes
        }]
    )
    return response.content[0].text

def extract_audio_from_video(file_bytes):
    with tempfile.NamedTemporaryFile(delete=False, suffix=".mp4") as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    clip = VideoFileClip(tmp_path)
    audio_path = tmp_path.replace(".mp4", ".mp3")
    clip.audio.write_audiofile(audio_path)
    clip.close()
    with open(audio_path, "rb") as f:
        return f.read()

# ─────────────────────────────────────
# PROCESAR ARCHIVO
# ─────────────────────────────────────
def process_file(uploaded_file):
    name = uploaded_file.name
    ext = name.split(".")[-1].lower()
    file_bytes = uploaded_file.read()

    if ext == "pdf":
        return extract_pdf(file_bytes)

    if ext == "docx":
        return extract_docx(file_bytes)

    if ext in ["txt", "md"]:
        return extract_txt(file_bytes)

    if ext == "pptx":
        return extract_pptx(file_bytes)

    if ext in ["mp3","wav","m4a"]:
        return transcribe_audio(file_bytes, name)

    if ext in ["mp4","mov"]:
        audio_bytes = extract_audio_from_video(file_bytes)
        return transcribe_audio(audio_bytes, name)

    return ""

# ─────────────────────────────────────
# GENERAR RESUMEN
# ─────────────────────────────────────
def generate_summary(text):
    r = client.messages.create(
        model="claude-sonnet-4-5-20250929",
        max_tokens=3000,
        messages=[{
            "role": "user",
            "content": f"Resume profesionalmente:\n{smart_truncate(text)}"
        }]
    )
    return r.content[0].text

# ─────────────────────────────────────
# CHAT
# ─────────────────────────────────────
def ask_question(question, context):
    r = client.messages.create(
        model="claude-sonnet-4-5-20250929",
        max_tokens=2000,
        system=f"Eres tutor experto. Contexto:\n{smart_truncate(context)}",
        messages=[{"role": "user", "content": question}]
    )
    return r.content[0].text

# ─────────────────────────────────────
# PRESENTACIÓN PRO
# ─────────────────────────────────────
def generate_presentation_data(text, num_slides):
    r = client.messages.create(
        model="claude-sonnet-4-5-20250929",
        max_tokens=min(9000, 600*num_slides),
        messages=[{
            "role": "user",
            "content": f"""
Crea una presentación profesional de {num_slides} diapositivas.
Máximo 5 bullets por slide.
Devuelve JSON:

{{
"titulo_general":"...",
"subtitulo":"...",
"color_primario":"1E2761",
"color_acento":"F96167",
"slides":[
  {{
    "titulo":"...",
    "contenido":["...","..."],
    "notas":"..."
  }}
]
}}

Contenido:
{smart_truncate(text)}
"""
        }]
    )
    return safe_json_parse(r.content[0].text)

def build_pptx(prs_data):
    prs = Presentation()

    primary = hex_to_rgb(prs_data.get("color_primario","1E2761"))
    accent = hex_to_rgb(prs_data.get("color_acento","F96167"))

    # PORTADA
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = prs_data.get("titulo_general","")
    slide.placeholders[1].text = prs_data.get("subtitulo","")

    # CONTENIDO
    for slide_data in prs_data.get("slides",[]):
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)

        slide.shapes.title.text = slide_data.get("titulo","")
        tf = slide.placeholders[1].text_frame
        tf.clear()

        for i,item in enumerate(slide_data.get("contenido",[])[:5]):
            if i==0:
                tf.text = item
            else:
                p = tf.add_paragraph()
                p.text = item
                p.level = 1

        for p in tf.paragraphs:
            p.font.size = Pt(20)
            p.font.color.rgb = accent

        if slide_data.get("notas"):
            slide.notes_slide.notes_text_frame.text = slide_data["notas"]

    bio = io.BytesIO()
    prs.save(bio)
    bio.seek(0)
    return bio.read()

# ─────────────────────────────────────
# STREAMLIT UI
# ─────────────────────────────────────
st.title("📚 AI Study Buddy")

if "text_data" not in st.session_state:
    st.session_state.text_data = ""
if "summary" not in st.session_state:
    st.session_state.summary = ""
if "chat" not in st.session_state:
    st.session_state.chat = []

uploaded_files = st.file_uploader(
    "Sube archivos",
    type=["pdf","docx","txt","md","pptx","mp3","wav","m4a","mp4","mov"],
    accept_multiple_files=True
)

if uploaded_files:
    combined = ""
    for f in uploaded_files:
        combined += process_file(f) + "\n\n"
    st.session_state.text_data = combined
    st.success("Archivos procesados.")

if st.button("Generar resumen"):
    st.session_state.summary = generate_summary(st.session_state.text_data)

if st.session_state.summary:
    st.markdown(st.session_state.summary)

st.divider()

# PPT
st.subheader("🎯 Generar presentación")
slides = st.slider("Número de diapositivas",6,25,10)

if st.button("Generar PPT"):
    prs_data = generate_presentation_data(st.session_state.text_data,slides)
    ppt_bytes = build_pptx(prs_data)
    st.download_button("Descargar PPT",
                       ppt_bytes,
                       "presentacion.pptx",
                       "application/vnd.openxmlformats-officedocument.presentationml.presentation")

st.divider()

# CHAT
st.subheader("💬 Preguntar")

question = st.text_input("Tu pregunta")
if st.button("Enviar"):
    answer = ask_question(question, st.session_state.text_data)
    st.session_state.chat.append(("Tú",question))
    st.session_state.chat.append(("Claude",answer))

for role,msg in st.session_state.chat:
    st.markdown(f"**{role}:** {msg}")
